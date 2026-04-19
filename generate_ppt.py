"""Polished 7-slide PPT for the Bypassing Detection project."""
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from copy import deepcopy

# ── Theme ──
NAVY   = RGBColor(0x0F, 0x2B, 0x4F)
ACCENT = RGBColor(0x1F, 0x87, 0xE0)
RED    = RGBColor(0xD9, 0x43, 0x51)
AMBER  = RGBColor(0xF5, 0xA6, 0x23)
GREEN  = RGBColor(0x2E, 0xA7, 0x5C)
GRAY   = RGBColor(0x5A, 0x6A, 0x7B)
LIGHT  = RGBColor(0xF4, 0xF7, 0xFA)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, x, y, w, h, fill=None, line_none=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill or WHITE
    if line_none: shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def textbox(slide, x, y, w, h, text, size=18, bold=False, color=None,
            align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color or NAVY
    r.font.name = "Calibri"
    return tb

def header(slide, title, subtitle=None):
    # Top bar
    rect(slide, 0, 0, SW, Inches(0.95), NAVY)
    # Accent strip
    rect(slide, 0, Inches(0.95), SW, Inches(0.05), ACCENT)
    textbox(slide, Inches(0.5), Inches(0.18), Inches(12), Inches(0.7),
            title, size=28, bold=True, color=WHITE)
    if subtitle:
        textbox(slide, Inches(0.5), Inches(1.15), Inches(12), Inches(0.4),
                subtitle, size=14, italic=True, color=GRAY)

def footer(slide, page, total=7):
    textbox(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
            "Bypassing Detection — Adversarial Attacks & Defenses",
            size=10, color=GRAY)
    textbox(slide, Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3),
            f"{page} / {total}", size=10, color=GRAY, align=PP_ALIGN.RIGHT)

def bullet_list(slide, x, y, w, h, items, size=16, gap=8):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            level, text, col = (it + (None,))[:3] if len(it) < 3 else it
        else:
            level, text, col = 0, it, None
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level; p.space_after = Pt(gap)
        bullet = "•  " if level == 0 else "–  "
        r = p.add_run(); r.text = bullet + text
        r.font.size = Pt(size if level == 0 else size - 2)
        r.font.bold = (level == 0 and col is None)
        r.font.color.rgb = col or (NAVY if level == 0 else GRAY)
        r.font.name = "Calibri"

def metric_card(slide, x, y, w, h, value, label, color=ACCENT):
    rect(slide, x, y, w, h, LIGHT)
    rect(slide, x, y, w, Inches(0.1), color)
    textbox(slide, x, y + Inches(0.3), w, Inches(0.9),
            value, size=34, bold=True, color=color, align=PP_ALIGN.CENTER)
    textbox(slide, x, y + Inches(1.2), w, Inches(0.4),
            label, size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════ SLIDE 1: TITLE ═══════════════
s = blank()
rect(s, 0, 0, SW, SH, NAVY)
# Decorative accent
rect(s, 0, Inches(3.4), SW, Inches(0.06), ACCENT)
textbox(s, Inches(1), Inches(2.2), Inches(11.3), Inches(1.2),
        "BYPASSING DETECTION", size=54, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(3.7), Inches(11.3), Inches(0.8),
        "Adversarial Attacks & Defenses on Road Sign Classifiers",
        size=24, color=RGBColor(180, 200, 220), align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(5.2), Inches(11.3), Inches(0.5),
        "4 Attack Methods  ·  5 Defense Strategies  ·  ResNet-34  ·  GTSRB Dataset",
        size=16, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ═══════════════ SLIDE 2: OVERVIEW ═══════════════
s = blank()
header(s, "Project Overview", "Evaluating model robustness under adversarial perturbations")

# Left: problem
textbox(s, Inches(0.5), Inches(1.7), Inches(6), Inches(0.5),
        "The Problem", size=20, bold=True, color=NAVY)
bullet_list(s, Inches(0.5), Inches(2.2), Inches(6), Inches(3),
    [
        "CNN classifiers power self-driving perception",
        "Imperceptible pixel perturbations cause misclassification",
        "A Stop sign can be read as Speed-limit",
        "Need: quantify how defenses hold up under attack",
    ], size=15)

# Right: approach
textbox(s, Inches(6.8), Inches(1.7), Inches(6), Inches(0.5),
        "Our Approach", size=20, bold=True, color=NAVY)
bullet_list(s, Inches(6.8), Inches(2.2), Inches(6), Inches(3),
    [
        "Train ResNet-34 on balanced GTSRB (4 classes)",
        "Implement 4 attacks: FGSM, PGD, GA, DE",
        "Implement 5 defenses across 3 categories",
        "Run 4×5 = 20 attack-defense combinations",
        "Build interactive web dashboard for demo",
    ], size=15)

# Bottom dataset strip
rect(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1), LIGHT)
textbox(s, Inches(0.7), Inches(5.9), Inches(12), Inches(0.4),
        "Dataset: GTSRB mapped 43 → 4 classes  |  Balanced 420/class train, 180/class test",
        size=14, bold=True, color=NAVY)
textbox(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.4),
        "Classes: crosswalk · speedlimit · stop · trafficlight",
        size=13, color=GRAY)
footer(s, 2)

# ═══════════════ SLIDE 3: METHODS ═══════════════
s = blank()
header(s, "Methods — Attacks & Defenses")

# Attacks table
textbox(s, Inches(0.5), Inches(1.3), Inches(6), Inches(0.4),
        "Attacks (4)", size=18, bold=True, color=RED)

attacks = [
    ("FGSM", "Gradient (ML)", "Single step, fast"),
    ("PGD", "Gradient (ML)", "Iterative, strong"),
    ("Genetic Algorithm", "Evolutionary", "Black-box, slow"),
    ("Differential Evolution", "Evolutionary", "Black-box, strongest"),
]
y = 1.85
rect(s, Inches(0.5), Inches(y), Inches(6), Inches(0.45), NAVY)
for i, col in enumerate(["Attack", "Type", "Property"]):
    textbox(s, Inches(0.5 + i * 2), Inches(y + 0.08), Inches(2), Inches(0.3),
            col, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for row_idx, row in enumerate(attacks):
    ry = y + 0.45 + row_idx * 0.48
    rect(s, Inches(0.5), Inches(ry), Inches(6), Inches(0.48),
         LIGHT if row_idx % 2 == 0 else WHITE)
    for i, val in enumerate(row):
        textbox(s, Inches(0.5 + i * 2), Inches(ry + 0.1), Inches(2), Inches(0.3),
                val, size=11, color=NAVY, align=PP_ALIGN.CENTER)

# Defenses table
textbox(s, Inches(6.9), Inches(1.3), Inches(6), Inches(0.4),
        "Defenses (5)", size=18, bold=True, color=GREEN)

defenses = [
    ("No Defense", "Baseline", "—"),
    ("Adversarial Training", "Proactive", "Train on PGD examples"),
    ("Defensive Distillation", "Proactive", "Soft labels T=20"),
    ("Input Transformation", "Pre-process", "Gaussian + median"),
    ("Detection Network", "Reactive", "Feature-based binary CNN"),
]
y = 1.85
rect(s, Inches(6.9), Inches(y), Inches(6.2), Inches(0.45), NAVY)
for i, col in enumerate(["Defense", "Category", "Approach"]):
    textbox(s, Inches(6.9 + i * 2.07), Inches(y + 0.08), Inches(2.07), Inches(0.3),
            col, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for row_idx, row in enumerate(defenses):
    ry = y + 0.45 + row_idx * 0.48
    rect(s, Inches(6.9), Inches(ry), Inches(6.2), Inches(0.48),
         LIGHT if row_idx % 2 == 0 else WHITE)
    for i, val in enumerate(row):
        textbox(s, Inches(6.9 + i * 2.07), Inches(ry + 0.1), Inches(2.07), Inches(0.3),
                val, size=10, color=NAVY, align=PP_ALIGN.CENTER)

# Bottom pipeline
rect(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.7), NAVY)
textbox(s, Inches(0.7), Inches(6.33), Inches(12), Inches(0.45),
        "Pipeline:  Image [0,1]  →  NormalizedModel  →  ResNet-34  →  Softmax  →  4 classes",
        size=14, bold=True, color=WHITE)
footer(s, 3)

# ═══════════════ SLIDE 4: TRAINING ═══════════════
s = blank()
header(s, "Training Pipeline & Model")

# Metric cards
metric_card(s, Inches(0.5), Inches(1.6), Inches(3), Inches(1.6),
            "100%", "Clean Val Accuracy", GREEN)
metric_card(s, Inches(3.7), Inches(1.6), Inches(3), Inches(1.6),
            "2,160", "Training Images", ACCENT)
metric_card(s, Inches(6.9), Inches(1.6), Inches(3), Inches(1.6),
            "ResNet-34", "Pretrained Backbone", NAVY)
metric_card(s, Inches(10.1), Inches(1.6), Inches(2.7), Inches(1.6),
            "H100", "GPU (Altair HPC)", AMBER)

# Training steps
textbox(s, Inches(0.5), Inches(3.5), Inches(12), Inches(0.4),
        "Training Stages", size=18, bold=True, color=NAVY)
steps = [
    ("1. Base Classifier", "10 epochs · AdamW · cosine LR · weighted CE"),
    ("2. Adversarial Training", "5 epochs · PGD ε=8/255 · 7 steps · α=2/255"),
    ("3. Defensive Distillation", "5 epochs · teacher → student · T=20 · α=0.7"),
    ("4. Detection Network", "5 epochs · 512-dim features · FGSM+PGD samples"),
    ("5. Full Evaluation", "4 attacks × 5 defenses = 20 combinations"),
]
for i, (lbl, detail) in enumerate(steps):
    y = 4.0 + i * 0.55
    rect(s, Inches(0.5), Inches(y), Inches(0.35), Inches(0.45), ACCENT)
    textbox(s, Inches(0.5), Inches(y + 0.05), Inches(0.35), Inches(0.4),
            str(i + 1), size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(s, Inches(1), Inches(y + 0.05), Inches(3.5), Inches(0.4),
            lbl, size=14, bold=True, color=NAVY)
    textbox(s, Inches(4.5), Inches(y + 0.05), Inches(8.5), Inches(0.4),
            detail, size=13, color=GRAY)
footer(s, 4)

# ═══════════════ SLIDE 5: RESULTS — ASR HEATMAP ═══════════════
s = blank()
header(s, "Results — Attack Success Rate (ASR)",
       "Lower is better  |  % of correctly-classified images fooled by attack")

# Load results
with open("results/evaluation_results.json") as f:
    data = json.load(f)

attacks_list = ["fgsm", "pgd", "genetic", "de"]
defenses_list = ["none", "adv_training", "input_transform", "detection", "distillation"]
attack_labels = ["FGSM", "PGD", "Genetic", "Diff. Evo."]
defense_labels = ["No Defense", "Adv Train", "Input Transform", "Detection", "Distillation"]

# Table dimensions
tx, ty = 1.0, 1.8
cell_w, cell_h = 2.15, 0.7
header_h = 0.6
row_header_w = 1.5

# Top-left cell
rect(s, Inches(tx), Inches(ty), Inches(row_header_w), Inches(header_h), NAVY)
textbox(s, Inches(tx), Inches(ty + 0.15), Inches(row_header_w), Inches(0.4),
        "Attack ↓ / Defense →", size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Defense headers
for j, dname in enumerate(defense_labels):
    x = tx + row_header_w + j * cell_w
    rect(s, Inches(x), Inches(ty), Inches(cell_w), Inches(header_h), NAVY)
    textbox(s, Inches(x), Inches(ty + 0.15), Inches(cell_w), Inches(0.4),
            dname, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Data rows
for i, akey in enumerate(attacks_list):
    y = ty + header_h + i * cell_h
    # Row label
    rect(s, Inches(tx), Inches(y), Inches(row_header_w), Inches(cell_h), NAVY)
    textbox(s, Inches(tx), Inches(y + 0.2), Inches(row_header_w), Inches(0.4),
            attack_labels[i], size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Cells
    for j, dkey in enumerate(defenses_list):
        x = tx + row_header_w + j * cell_w
        asr = data["results"][akey][dkey].get("attack_success_rate", 0)
        # Color scale
        if asr < 10:
            fill = GREEN
        elif asr < 40:
            fill = RGBColor(0x8A, 0xC9, 0x26)
        elif asr < 70:
            fill = AMBER
        elif asr < 90:
            fill = RGBColor(0xE5, 0x70, 0x2B)
        else:
            fill = RED
        rect(s, Inches(x), Inches(y), Inches(cell_w), Inches(cell_h), fill)
        textbox(s, Inches(x), Inches(y + 0.2), Inches(cell_w), Inches(0.4),
                f"{asr:.1f}%", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Legend
ly = ty + header_h + 4 * cell_h + 0.3
textbox(s, Inches(tx), Inches(ly), Inches(2), Inches(0.3),
        "Color scale:", size=11, bold=True, color=NAVY)
legend_items = [("< 10%", GREEN), ("10-40%", RGBColor(0x8A, 0xC9, 0x26)),
                ("40-70%", AMBER), ("70-90%", RGBColor(0xE5, 0x70, 0x2B)),
                ("> 90%", RED)]
lx = tx + 1.3
for label, clr in legend_items:
    rect(s, Inches(lx), Inches(ly), Inches(0.35), Inches(0.3), clr)
    textbox(s, Inches(lx + 0.4), Inches(ly), Inches(1), Inches(0.3),
            label, size=10, color=GRAY)
    lx += 1.5
footer(s, 5)

# ═══════════════ SLIDE 6: FINDINGS ═══════════════
s = blank()
header(s, "Key Findings & Insights")

# Left column - winners
rect(s, Inches(0.5), Inches(1.5), Inches(6), Inches(2.5), LIGHT)
rect(s, Inches(0.5), Inches(1.5), Inches(6), Inches(0.1), GREEN)
textbox(s, Inches(0.7), Inches(1.7), Inches(6), Inches(0.5),
        "Strongest Defenses", size=18, bold=True, color=GREEN)
bullet_list(s, Inches(0.7), Inches(2.25), Inches(5.6), Inches(2),
    [
        "Distillation: 0.8% ASR vs PGD (best proactive)",
        "Adversarial Training: 9.4% vs FGSM, 20.5% vs PGD",
        "Detection Network: 0% vs FGSM (catches gradient attacks)",
    ], size=13, gap=6)

# Right column - losers
rect(s, Inches(6.8), Inches(1.5), Inches(6), Inches(2.5), LIGHT)
rect(s, Inches(6.8), Inches(1.5), Inches(6), Inches(0.1), RED)
textbox(s, Inches(7), Inches(1.7), Inches(6), Inches(0.5),
        "Weakest Combos", size=18, bold=True, color=RED)
bullet_list(s, Inches(7), Inches(2.25), Inches(5.6), Inches(2),
    [
        "DE vs Distillation: 97.9% ASR (bypasses soft labels)",
        "Genetic vs Distillation: 85.8% (EC evades smoothing)",
        "Input Transform vs PGD: 53.8% (fails at high ε)",
    ], size=13, gap=6)

# Bottom takeaways
rect(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.5), NAVY)
textbox(s, Inches(0.7), Inches(4.45), Inches(12), Inches(0.5),
        "Takeaways", size=18, bold=True, color=WHITE)
takeaways = [
    "Gradient defenses (adv training, distillation) block gradient attacks but fail against black-box EC",
    "Evolutionary attacks (Genetic, DE) bypass defenses because they don't rely on gradients",
    "Detection works well against ML attacks but is blind to evolutionary perturbations",
    "No single defense is universal — layered defenses are essential for robustness",
]
for i, t in enumerate(takeaways):
    textbox(s, Inches(0.9), Inches(4.95 + i * 0.42), Inches(12), Inches(0.4),
            f"→  {t}", size=13, color=WHITE)
footer(s, 6)

# ═══════════════ SLIDE 7: DEMO & CONCLUSION ═══════════════
s = blank()
header(s, "Demo, Stack & Conclusion")

# Demo panel
rect(s, Inches(0.5), Inches(1.4), Inches(6), Inches(3), LIGHT)
rect(s, Inches(0.5), Inches(1.4), Inches(0.1), Inches(3), ACCENT)
textbox(s, Inches(0.8), Inches(1.55), Inches(6), Inches(0.5),
        "Live Demo — Interactive Web Lab", size=17, bold=True, color=ACCENT)
bullet_list(s, Inches(0.8), Inches(2.1), Inches(5.5), Inches(2.5),
    [
        "Upload/select road sign image",
        "Pick attack + epsilon slider (0-0.15)",
        "See perturbed image, prediction shift, confidence bars",
        "Compare all 5 defense results side-by-side",
    ], size=13, gap=5)

# Stack panel
rect(s, Inches(6.8), Inches(1.4), Inches(6), Inches(3), LIGHT)
rect(s, Inches(6.8), Inches(1.4), Inches(0.1), Inches(3), AMBER)
textbox(s, Inches(7.1), Inches(1.55), Inches(6), Inches(0.5),
        "Tech Stack", size=17, bold=True, color=AMBER)
bullet_list(s, Inches(7.1), Inches(2.1), Inches(5.5), Inches(2.5),
    [
        "PyTorch 2.6 (CUDA 12.1) · torchvision",
        "Flask + HTML/CSS/JS dashboard",
        "NumPy · matplotlib · PIL",
        "Altair HPC / pytorch_pbs Docker image",
    ], size=13, gap=5)

# Conclusion bar
rect(s, Inches(0.5), Inches(4.65), Inches(12.3), Inches(2.15), NAVY)
textbox(s, Inches(0.7), Inches(4.8), Inches(12), Inches(0.5),
        "Conclusion", size=18, bold=True, color=WHITE)
concl = [
    "Successfully implemented and evaluated 20 attack-defense combinations",
    "Achieved 100% clean accuracy; robust accuracy varies 0-100% depending on pairing",
    "Future: ensemble defenses, certified robustness, adaptive attacks, larger dataset",
]
for i, t in enumerate(concl):
    textbox(s, Inches(0.9), Inches(5.3 + i * 0.42), Inches(12), Inches(0.4),
            f"✓  {t}", size=13, color=WHITE)
footer(s, 7)

prs.save("Bypassing_Detection_Presentation.pptx")
print("Saved: Bypassing_Detection_Presentation.pptx (7 slides)")
